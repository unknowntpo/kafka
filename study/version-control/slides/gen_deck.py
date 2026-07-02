#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
import os

HERE = os.path.dirname(os.path.abspath(__file__))
ICO = os.path.join(HERE, "icons") + os.sep

ACCENT = RGBColor(0x18, 0x5F, 0xA5)
DARK   = RGBColor(0x1E, 0x20, 0x24)
GRAY   = RGBColor(0x5f, 0x5e, 0x5a)
MUTED  = RGBColor(0x8a, 0x8a, 0x8a)
CODEBG = RGBColor(0xF3, 0xF1, 0xEC)   # subtle warm-grey panel (Claude minimalist)
CODEFG = RGBColor(0x2C, 0x2C, 0x2A)
HAIR   = RGBColor(0xE3, 0xE0, 0xD8)   # hairline border
SOLVEBG= RGBColor(0xE9, 0xF3, 0xDE)   # soft green tint
SOLVEFG= RGBColor(0x27, 0x50, 0x0A)
GREEN  = RGBColor(0x27, 0x50, 0x0A)

UP = "b7b1c0a83d856766390ee0c05e33b63711eee80e"   # apache/kafka trunk commit (upstream line numbers)
FK = "27342708500d978cb9b7477b4a870d4954ec13b1"   # unknowntpo/kafka commit (demo tests / lab)
P = {
 "protocol.md":"docs/design/protocol.md",
 "zk2kraft.md":"docs/getting-started/zk2kraft.md",
 "upgrade.md":"docs/getting-started/upgrade.md",
 "NodeApiVersions.java":"clients/src/main/java/org/apache/kafka/clients/NodeApiVersions.java",
 "NetworkClient.java":"clients/src/main/java/org/apache/kafka/clients/NetworkClient.java",
 "ApiKeys.java":"clients/src/main/java/org/apache/kafka/common/protocol/ApiKeys.java",
 "FetchRequest.json":"clients/src/main/resources/common/message/FetchRequest.json",
 "FetchRequest.java":"clients/src/main/java/org/apache/kafka/common/requests/FetchRequest.java",
 "RequestContext.java":"clients/src/main/java/org/apache/kafka/common/requests/RequestContext.java",
 "SocketServer.scala":"core/src/main/scala/kafka/network/SocketServer.scala",
 "KafkaApis.scala":"core/src/main/scala/kafka/server/KafkaApis.scala",
 "TransactionMarkerChannelManager.scala":"core/src/main/scala/kafka/coordinator/transaction/TransactionMarkerChannelManager.scala",
 "OffsetsForLeaderEpochRequest.java":"clients/src/main/java/org/apache/kafka/common/requests/OffsetsForLeaderEpochRequest.java",
 "RequestResponseTest.java":"clients/src/test/java/org/apache/kafka/common/requests/RequestResponseTest.java",
 "MessageTest.java":"clients/src/test/java/org/apache/kafka/common/message/MessageTest.java",
 "Dockerfile":"tests/docker/Dockerfile",
 "version.py":"tests/kafkatest/version.py",
 "client_compat.py":"tests/kafkatest/tests/client/client_compatibility_features_test.py",
 "upgrade_test.py":"tests/kafkatest/tests/core/upgrade_test.py",
 "MetadataVersionTest.java":"server-common/src/test/java/org/apache/kafka/server/common/MetadataVersionTest.java",
 "BrokerBlockingSender.scala":"core/src/main/scala/kafka/server/BrokerBlockingSender.scala",
 "MetadataVersion.java":"server-common/src/main/java/org/apache/kafka/server/common/MetadataVersion.java",
 "RemoteLeaderEndPoint.scala":"core/src/main/scala/kafka/server/RemoteLeaderEndPoint.scala",
 "FeatureControlManager.java":"metadata/src/main/java/org/apache/kafka/controller/FeatureControlManager.java",
 "ClusterControlManager.java":"metadata/src/main/java/org/apache/kafka/controller/ClusterControlManager.java",
 "Feature.java":"server-common/src/main/java/org/apache/kafka/server/common/Feature.java",
 "NodeApiVersionsTest.java":"clients/src/test/java/org/apache/kafka/clients/NodeApiVersionsTest.java",
 "ClusterControlManagerTest.java":"metadata/src/test/java/org/apache/kafka/controller/ClusterControlManagerTest.java",
 "docker-compose.yml":"study/version-control/lab/tour0/docker-compose.yml",
}
KIP896="https://cwiki.apache.org/confluence/display/KAFKA/KIP-896%3A+Remove+old+client+protocol+API+versions+in+Kafka+4.0"
KIP35="https://cwiki.apache.org/confluence/display/KAFKA/KIP-35+-+Retrieving+protocol+version"
KIP584="https://cwiki.apache.org/confluence/display/KAFKA/KIP-584%3A+Versioning+scheme+for+features"
PGURL="https://www.postgresql.org/docs/current/protocol-flow.html"
MONGO="https://www.mongodb.com/docs/manual/reference/command/hello/"

def A(name, line=None, disp=None):
    u = "https://github.com/apache/kafka/blob/%s/%s" % (UP, P[name])
    if line: u += "#L%d" % line
    return (disp or (name + (":%d" % line if line else "")), u, True)
def F(name, line=None, disp=None):
    u = "https://github.com/unknowntpo/kafka/blob/%s/%s" % (FK, P[name])
    if line: u += "#L%d" % line
    return (disp or (name + (":%d" % line if line else "")), u, True)
def X(text, url): return (text, url, False)
def T(text): return (text, None, False)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
CW = 12.0

def tb(slide, x, y, w, h):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    return b, tf

def run(p, text, size, color, bold=False, mono=False):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    r.font.name = "Courier New" if mono else "Calibri"
    return r

def has_cjk(s):
    return any('　' <= c <= '鿿' or '＀' <= c <= '￯' for c in s)

def est_lines(text, cpl):
    n = 0
    for line in text.split("\n"):
        n += max(1, -(-len(line)//cpl))
    return n

def render_ref(s, tokens):
    _, tf = tb(s, 0.7, 7.0, CW, 0.42)
    p = tf.paragraphs[0]
    run(p, "REF   ", 12, MUTED, bold=True)
    for i, (text, url, mono) in enumerate(tokens):
        if i > 0: run(p, "   ·   ", 12, MUTED)
        r = run(p, text, 12, ACCENT if url else MUTED, mono=mono)
        if url:
            r.hyperlink.address = url
            r.font.underline = True

def add_content(kicker, title, blocks, ref=None):
    s = prs.slides.add_slide(BLANK)
    if kicker:
        _, tf = tb(s, 0.7, 0.5, CW, 0.4); run(tf.paragraphs[0], kicker, 13, ACCENT, bold=True)
    _, tf = tb(s, 0.7, 0.92, CW, 1.1); run(tf.paragraphs[0], title, 32, DARK, bold=True)
    y = 2.35
    for blk in blocks:
        kind, text = blk[0], blk[1]
        if kind == "code":
            h = est_lines(text, 58)*0.30 + 0.30
            box, tf = tb(s, 0.7, y, CW, h)
            box.fill.solid(); box.fill.fore_color.rgb = CODEBG
            box.line.color.rgb = HAIR; box.line.width = Pt(0.5)
            tf.margin_left = Inches(0.16); tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.1)
            mono = not has_cjk(text)
            first = True
            for ln in text.split("\n"):
                p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
                run(p, ln if ln else " ", 13.5, CODEFG, mono=mono)
            y += h + 0.22
        elif kind == "solve":
            h = est_lines(text, 50)*0.34 + 0.28
            box, tf = tb(s, 0.7, y, CW, h)
            box.fill.solid(); box.fill.fore_color.rgb = SOLVEBG; box.line.fill.background()
            tf.margin_left = Inches(0.16); tf.margin_top = Inches(0.09)
            run(tf.paragraphs[0], text, 15, SOLVEFG)
            y += h + 0.20
        elif kind == "note":
            h = est_lines(text, 62)*0.32 + 0.16
            _, tf = tb(s, 0.7, y, CW, h); run(tf.paragraphs[0], text, 15, GRAY)
            y += h + 0.14
        else:
            icon = blk[2] if len(blk) > 2 else "point"
            h = est_lines(text, 42)*0.38 + 0.16
            s.shapes.add_picture(ICO + icon + ".png", Inches(0.72), Inches(y + 0.03), height=Inches(0.30))
            _, tf = tb(s, 1.18, y, CW - 0.5, h); run(tf.paragraphs[0], text, 16.5, DARK)
            y += h + 0.16
    if ref: render_ref(s, ref)
    return s

def add_quiz(tag, q, opts, correct, why, ref, reveal):
    # reveal=False → 只出題、蓋住答案；reveal=True → 公布正解（下一頁）
    s = prs.slides.add_slide(BLANK)
    label = "隨堂考 · " + tag + ("　（正解）" if reveal else "")
    _, tf = tb(s, 0.7, 0.5, CW, 0.4); run(tf.paragraphs[0], label, 13, ACCENT, bold=True)
    _, tf = tb(s, 0.7, 0.92, CW, 1.3); run(tf.paragraphs[0], q, 26, DARK, bold=True)
    y = 2.55
    letters = "ABCD"
    for k, o in enumerate(opts):
        h = est_lines(o, 52)*0.36 + 0.18
        box, tf = tb(s, 0.9, y, 11.5, h)
        ok = reveal and (k == correct)
        if ok:
            box.fill.solid(); box.fill.fore_color.rgb = SOLVEBG; box.line.fill.background()
            tf.margin_left = Inches(0.14); tf.margin_top = Inches(0.06)
        p = tf.paragraphs[0]
        run(p, letters[k] + ".  ", 16, GREEN if ok else DARK, bold=ok)
        run(p, o + ("   ✓" if ok else ""), 16, GREEN if ok else DARK, bold=ok)
        y += h + 0.14
    _, tf = tb(s, 0.7, y + 0.1, CW, 1.0)
    if reveal:
        p = tf.paragraphs[0]; run(p, "正解 " + letters[correct] + " — ", 14, GREEN, bold=True); run(p, why, 14, DARK)
    else:
        run(tf.paragraphs[0], "想一想 —— 下一頁公布答案", 14, MUTED)
    if ref and reveal: render_ref(s, ref)
    return s
def quiz_pair(i):
    add_quiz(*QUIZ[i], reveal=False)
    add_quiz(*QUIZ[i], reveal=True)

# 只留本場範圍內、有教育意義的 2 題（其餘屬兄弟場或過於簡單，已移除）
QUIZ = [
 ("核心", "replica fetch（broker↔broker）的 Fetch 版本怎麼決定？",
  ["兩 broker ApiVersions 協商取交集", "由 finalized metadata.version 決定", "永遠用最新版", "controller 每次指定"], 1,
  "fetchRequestVersion(MV)：MV 是集中共識，不做 per-connection 協商。",
  [A("RemoteLeaderEndPoint.scala",215), A("MetadataVersion.java",273)]),
 ("核心", "client 繞過協商、直接送出 broker 不支援的 API version，會怎樣？",
  ["一律回 UNSUPPORTED_VERSION 錯誤碼", "只有 ApiVersions 例外回錯誤碼+支援範圍；其他 API 直接關閉連線", "broker 自動降版處理", "忽略版本照常處理"], 1,
  "ApiVersions 是協商的第一支 request（此時雙方還不知道彼此版本），若它也關閉連線，client 永遠問不到支援範圍，故回 v0+錯誤碼讓 client 重新協商；其他 API 在 RequestContext 解析失敗、SocketServer 關閉連線。",
  [A("RequestContext.java",112), A("SocketServer.scala",781)]),
]
# ---- Title ----
s = prs.slides.add_slide(BLANK)
bar = s.shapes.add_shape(1, Inches(0.9), Inches(2.62), Inches(0.1), Inches(1.5))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
_, tf = tb(s, 1.2, 1.5, 11.3, 0.5); run(tf.paragraphs[0], "幾萬個節點之上的版本控制 · 第三場", 15, ACCENT, bold=True)
_, tf = tb(s, 1.2, 2.6, 11.3, 1.4); run(tf.paragraphs[0], "溝通當下的版本選擇", 46, DARK, bold=True)
_, tf = tb(s, 1.2, 4.2, 11.3, 0.9); run(tf.paragraphs[0], "長壽的 client、滾動升級的 broker，怎麼喬出共同版本", 20, GRAY)
_, tf = tb(s, 1.2, 6.65, 11.3, 0.5); run(tf.paragraphs[0], "Eric Chang · 2026", 14, MUTED)

add_content("議程", "本場範圍", [
    ("bullet", "Part 1 — 為什麼不能「一個版本打天下」，版本又是怎麼選出來的", "point"),
    ("bullet", "Part 2 — 協商失敗時，會看到什麼訊息", "point"),
])

add_content("點題", "為什麼不能「一個版本打天下」？", [
    ("code", "直覺（想像中）：  client v4.1 ─────── broker v4.1   同版、一起升"),
    ("bullet", "現實一：client 是內嵌在應用裡的函式庫——IoT 裝置、車載系統、多年前的 batch job，隨應用長期運行、難更新（Kafka 曾為此保留每個 protocol 版本近九年）", "clock-hour-4"),
    ("bullet", "現實二：Kafka 要求 zero downtime → broker 只能一台一台滾動升級 → 過程必然新舊混版", "server-2"),
    ("solve", "兩個現實加起來：不能鎖全叢集同一版（又要不停機）——版本無法事先統一，那要怎麼選？"),
    ("note", "這其實就是 backward／forward compatibility：一般服務多靠『加欄位＋ URL 版號 /v1 /v2』相容；Kafka 特別在每支 API 各自宣告版本區間、連線當下協商挑一版（更像 TLS 選 cipher，而非固定版號）。"),
], [X("KIP-896",KIP896), A("protocol.md",94)])

add_content("Part 1 · 術語", "要談「怎麼選」，先分清楚「版本」是哪一層", [
    ("bullet", "上一頁的結論是「版本無法事先統一」；但在講機制前得先知道，日常講的『版本』其實混了三個東西", "arrows-split"),
    ("code", "release version   = 我裝了哪版 binary            (per-node)\nmetadata.version  = 叢集 finalized 的 feature level (cluster-wide)\nwire API version  = 這條連線講第幾版               (per-connection)"),
    ("note", "finalize＝管理員手動宣告全叢集一致採用的 feature level（第一場《版本定義》的主題）。本場主角是 wire 版本。"),
], [T("kafka-features describe"), A("zk2kraft.md",71)])

add_content("Part 1 · 架構 (a)", "版本怎麼定：協商，還是由 MV 事先釘版", [
    ("bullet", "機制一 · 協商（KIP-35）：連線後送 ApiVersionsRequest、取交集最高——多數路徑都是", "arrows-split"),
    ("bullet", "機制二 · 由叢集 MV 事先釘版（KIP-584）：只有 replica fetch（partition replication）這條不協商", "ban"),
    ("code", "協商定版   client↔broker · broker↔controller · Raft · broker↔broker txn markers\nMV 釘版    replica fetch（follower → leader partition replication）—— 唯一一條"),
    ("note", "所以「broker↔broker 不協商」並不精確：只有 replica fetch 這條；同是 broker↔broker 的 transaction markers 仍協商。"),
], [X("KIP-35",KIP35), X("KIP-584",KIP584), A("BrokerBlockingSender.scala",95,"replica fetch=false"), A("TransactionMarkerChannelManager.scala",99,"txn markers=true")])

add_content("Part 1 · 架構 (b)", "查詢之後，最終版本誰說了算？", [
    ("code", "協商           取交集最高（client↔broker、broker↔controller、Raft、txn markers）\nreplica fetch  不協商——Fetch/ListOffsets 由 finalized MV、OffsetsForLeaderEpoch 寫死 v4"),
    ("bullet", "replica fetch 連 ApiVersionsRequest 都不送（discoverBrokerVersions=false）：這條路上一整組 RPC 版本都不看對端能力，改由叢集事先決定", "ban"),
    ("solve", "為什麼 replica fetch 交給 MV？它要求 follower↔leader 全講同一版、沒法各連線各挑。代價『升級變兩階段』其實是 MV 治理的通性（凡版本由 MV 決定者皆然），非它專屬；協商連線則自動挑上去"),
], [A("MetadataVersion.java",273,"fetchRequestVersion"), A("BrokerBlockingSender.scala",95,"discoverBrokerVersions=false"), A("OffsetsForLeaderEpochRequest.java",65,"forFollower 寫死 v4"), A("RemoteLeaderEndPoint.scala",215)])

# §2a — 數線圖：架構下的一個舉例（cb 協商 vs bb 由 MV）
s2a = prs.slides.add_slide(BLANK)
_, tf = tb(s2a, 0.7, 0.5, CW, 0.4); run(tf.paragraphs[0], "Part 1 · 舉例", 13, ACCENT, bold=True)
_, tf = tb(s2a, 0.7, 0.9, CW, 0.7); run(tf.paragraphs[0], "以一支 Fetch 為例：consumer 協商、replica 由 MV", 29, DARK, bold=True)
_, tf = tb(s2a, 0.7, 1.62, CW, 0.4); run(tf.paragraphs[0], "同一支 Fetch 兩種身分：consumer fetch 走 client↔broker、replica fetch 走 broker↔broker", 14, GRAY)
s2a.shapes.add_picture(ICO + "s2a-rangeline.png", Inches(1.82), Inches(2.15), width=Inches(9.7))
render_ref(s2a, [A("FetchRequest.java",165,"forConsumer/forReplica"), A("MetadataVersion.java",273,"fetchRequestVersion"), A("FetchRequest.json",61,"validVersions")])

add_content("Part 1 · 對照", "同一道題，別家怎麼答——各付什麼代價", [
    ("bullet", "Kafka：per-API 協商，replica fetch 交給 MV 釘版。買到單支 API 獨立演進、複製層可線上滾動升級；代價＝協定面積最大、相容性測試按 API × 版本數放大", "arrows-split"),
    ("bullet", "MongoDB：全域一個 wire version＋FCV 治理叢集。實作簡單；代價＝粒度粗、無法單支 API 獨立演進（FCV 與 MV 同一套思路）", "database"),
    ("bullet", "PostgreSQL：協定 v3 逾二十年極穩定；但實體複製鎖 major → 複製層無法線上滾動升級（得停機 pg_upgrade 或另建叢集）", "leaf"),
    ("solve", "沒有免費的選擇：Kafka 拿協定複雜度換到細粒度演進＋複製層線上升級——後者正是 PostgreSQL 買不到的"),
], [X("MongoDB hello / FCV",MONGO), X("PostgreSQL protocol",PGURL)])

add_content("Part 1 · 代價", "Kafka 這個選擇有多貴（實測數字）", [
    ("code", "90 支 API × 308 個現役 wire 版本 → generator 吹成約 18× 生成碼（~18 萬行）\n每加一版：改 .json、逐版生 read/write/size、ApiKeys 補版號"),
    ("bullet", "broker handler 還得逐版分岔：按協商版本讀/填欄位、回應用『同一版』序列化、老版缺的欄位給預設，連語意差異（如 topic-id 取代 topic-name）都要 if (version≥N)", "arrows-split"),
    ("note", "這還只是「生成碼＋handler」；相容性測試是另一條隱形帳單（下兩張細看）——build 全版本 round-trip＋release 養一堆歷史版本。"),
    ("solve", "貴到 Kafka 自己在 KIP-896 承認「maintenance cost up, value down」，4.0 砍掉 2.1 以前的舊版本止血"),
], [X("KIP-896",KIP896), A("RequestContext.java",137,"buildResponseSend"), A("KafkaApis.scala",568,"fetch version()>=13"), A("FetchRequest.json",80,"MaxBytes v3+ default")])
add_content("Part 1 · 代價 · 測試", "相容性測試（一）：每次 build 全版本 round-trip", [
    ("code", "RequestResponseTest.testSerialization：ApiKeys × allVersions 全積\n94 支 RPC、308 個現役請求版本 → build → 序列化 → 反序列化 → 比對，逐一跑"),
    ("bullet", "靠 ApiKeys.values() × allVersions() 迴圈，新增版本自動涵蓋——省了手改，代價是每次 build 執行量只增不減", "arrows-split"),
    ("bullet", "想跳過某個歷史版本？得進 toSkip 白名單特案處理——目前只 4 支 RPC 拿得到豁免", "ban"),
    ("solve", "KIP-896 直言維護舊版本「both in code complexity and the testing matrix」——測試面正是其一"),
], [A("RequestResponseTest.java",340,"testSerialization"), A("MessageTest.java",716,"round-trip"), X("KIP-896",KIP896)])

add_content("Part 1 · 代價 · 測試", "相容性測試（二）：把過去七年的 Kafka 一起養著", [
    ("code", "系統測試映像 curl 下載 22 個歷史版本二進位（2.1.1 → 4.3.0）\nclient 相容測試：每個歷史 broker 各跑一遍（22 + dev = 23 param × 2 支測試）"),
    ("bullet", "升級/降級是 from×to 矩陣：4 個 @matrix × 11 個 from_version；混版交易再 22 組；MetadataVersion 另有 25 個現役 MV × 7 個 enum 測試", "arrows-split"),
    ("solve", "取捨：相容承諾＝「一次寫、永遠不能少測」——每發一個 release，version.py＋Dockerfile＋各測試 @parametrize 都得手動 +1"),
], [A("Dockerfile",79,"下載 22 版"), A("client_compat.py",109,"client 相容 ×23"), A("upgrade_test.py",167,"升級矩陣"), A("MetadataVersionTest.java",219,"25 MV")])

quiz_pair(0)  # 小測驗 1：replica fetch 版本誰決定

# ---- Part 2：失敗會有什麼訊息 ----
add_content("Part 2 · 失敗訊息", "通訊當下協商不出版本，會看到什麼", [
    ("code", "交集空                       → client 端 UnsupportedVersionException（送出前 abort）\nclient 繞過協商、硬送不支援版本 → broker 丟 UnsupportedVersionException、關閉連線\n（ApiVersions 例外：回 v0 + UNSUPPORTED_VERSION 錯誤碼）"),
    ("bullet", "多數情況 client 在送出前就本地中止，根本沒上網路", "ban"),
    ("note", "「硬送」指 client↔broker：自刻或有 bug 的 client 繞過協商；broker↔broker 版本已由 MV 事先釘住，不會送出對方不懂的版本。"),
], [A("NodeApiVersions.java",149,"latestUsableVersion"), A("NetworkClient.java",591,"NetworkClient"), A("RequestContext.java",112,"RequestContext")])

add_content("Part 2 · 失敗訊息", "版本截斷：為什麼「升一點點」不夠", [
    ("bullet", "前一頁「交集空」最容易被忽略的根因——舊 wire 版本被整批移除", "alert-triangle"),
    ("code", "舊 wire API version 被整個移除，min 可 > 0\n  Fetch：4–18（v0–3 在 4.0 移除）"),
    ("solve", "解法：升任一端到 4.0 前，先確認另一端 ≥ 2.1（upgrade guide 明訂雙向要求），否則協商結果沒有可用版本"),
    ("note", "取捨：0.8.0（2013）起九年保留每個版本——每個舊版本都是活的程式碼與測試面；4.0 用「斷 2018 年前的 client」換「刪九年的碼」，只有 major 版本邊界付得起。"),
    ("note", "銜接（非本場）：finalize / 升降當下的錯誤（INVALID_UPDATE_VERSION 等）屬「運行時的版本升降」那場，本場不展開。"),
], [A("upgrade.md",229), A("FetchRequest.json",61,"FetchRequest.json validVersions")])
quiz_pair(1)  # 小測驗 2：自刻不支援版本會怎樣

add_content("Recap", "回到開場那個問題", [
    ("code", "想像中：client v4.1 ─── broker v4.1\n實際上：同一顆 4.1 broker，同時講 Fetch v11（對老 consumer）和 v17（對 replica）"),
    ("bullet", "Part 1：不能「一個版本打天下」（client 長壽、broker 滾動升級）；多數路徑靠協商，唯一不協商、由 finalized MV 釘版的是 replica fetch（其餘 broker↔broker 仍協商）", "point"),
    ("bullet", "Part 2：協商不出版本 → UnsupportedVersionException（多在送出前中止）；finalize／升降的錯誤交給運行時那場", "point"),
    ("solve", "立場：per-API 細粒度協商＋replication 集中治理，是為「client 生態極度分散」量身的取捨——對 Kafka 划算，但不是通用解"),
])

out = os.path.join(HERE, "..", "kafka-version-negotiation.pptx")
prs.save(out)
print("slides:", len(prs.slides._sldIdLst), "->", out)
